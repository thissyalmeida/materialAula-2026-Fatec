from pptx import Presentation
import os
import re
import uuid

# ========================================
# CONFIGURAÇÕES GERAIS
# ========================================

PASTA_ENTRADA = "slides"
PASTA_SAIDA = "aulas"
PASTA_IMAGENS = "assets/img"

PALAVRAS_QUIZ = [
    "pratique",
    "exercício",
    "fixação",
    "teste"
]

PALAVRAS_EXERCICIO = [
    "pratique",
    "exercício"
]

EMOJIS = {
    "lógica": "🧠",
    "algoritmo": "⚙️",
    "programa": "💻",
    "condicional": "🔀",
    "repetição": "🔁",
    "função": "🧩",
    "vetor": "📦",
    "matriz": "📊",
    "entrada": "📥",
    "saída": "📤"
}

# ========================================
# UTILIDADES
# ========================================

def emoji_para(texto):
    texto = texto.lower()

    for palavra, emoji in EMOJIS.items():
        if palavra in texto:
            return emoji

    return "📌"


def limpar_texto(texto):
    texto = texto.strip()
    texto = re.sub(r'\s+', ' ', texto)
    return texto


# ========================================
# DETECÇÃO DE CÓDIGO C
# ========================================

def parece_codigo_c(texto):

    padroes = [
        r"#include",
        r"int main",
        r"printf",
        r"scanf",
        r"\{",
        r"\}"
    ]

    for p in padroes:
        if re.search(p, texto):
            return True

    return False


def caixa_codigo(codigo):

    return f"""```c{codigo}"""

#========================================
#EXTRAÇÃO DE IMAGENS
#========================================

def extrair_imagem(shape, slide_num):

    if not os.path.exists(PASTA_IMAGENS):
        os.makedirs(PASTA_IMAGENS)

    image = shape.image
    image_bytes = image.blob

    nome = f"slide_{slide_num}_{uuid.uuid4().hex[:6]}.png"

    caminho = os.path.join(
        PASTA_IMAGENS,
        nome
    )

    with open(caminho, "wb") as f:
        f.write(image_bytes)

    return caminho
#========================================
#COMPONENTES VISUAIS
#========================================

def gerar_header(nome):

    return f"""# 🚀 {nome} 🎯 Objetivo da aula:Aprender os conceitos fundamentais desta aula."""

def caixa_exercicio_html(titulo, conteudo):

    return f""" 🧪 {titulo} <div class="exercise-box"> {conteudo} <br> <input type="text" placeholder="Digite sua resposta aqui"> <br><br> <button onclick="alert('Resposta registrada!')"> ✔️ Verificar resposta </button> </div>"""

def criar_quiz_html(pergunta, alternativas):

    html = f"""🧪 {pergunta}<div class="quiz-box"> <form>"""

    for alt in alternativas:

        html += f"""<label> <input type="radio" name="quiz"> {alt} </label><br> """
    html += """<br><button type="button" onclick="alert('Resposta registrada!')">✔️ Verificar</button> </form> </div>"""

    return html

def eh_quiz(texto):

    texto = texto.lower()

    return any(
        palavra in texto
        for palavra in PALAVRAS_QUIZ
    )

def eh_exercicio(texto):

    texto = texto.lower()

    return any(
        palavra in texto
        for palavra in PALAVRAS_EXERCICIO
    )
#========================================
#PROCESSAMENTO PRINCIPAL
#========================================

def converter_pptx(arquivo):

    prs = Presentation(arquivo)

    nome = os.path.splitext(
        os.path.basename(arquivo)
    )[0]

    markdown = []

    markdown.append(
        gerar_header(nome)
    )

    for slide_num, slide in enumerate(prs.slides):

        textos = []

        for shape in slide.shapes:

            # TEXTO
            if hasattr(shape, "text"):

                texto = limpar_texto(
                shape.text
                )

                if texto:
                    textos.append(texto)

            # IMAGEM
            if shape.shape_type == 13:

                caminho = extrair_imagem(
                    shape,
                    slide_num
                )

                markdown.append(
                    f"\n![Imagem](/{caminho})\n"
                )

        if not textos:
            continue

        titulo = textos[0]

        emoji = emoji_para(titulo)

        # QUIZ
        if eh_quiz(titulo):

            alternativas = textos[1:]

            markdown.append(
                criar_quiz_html(
                    titulo,
                    alternativas
                )
            )

            continue

        # EXERCÍCIO
        if eh_exercicio(titulo):

            conteudo = "\n".join(
                [f"- {t}" for t in textos[1:]]
            )

            markdown.append(
                caixa_exercicio_html(
                    titulo,
                    conteudo
                )
            )

            continue

        # SEÇÃO NORMAL
        markdown.append(
            f"\n## {emoji} {titulo}\n"
        )

        for linha in textos[1:]:

            if parece_codigo_c(linha):

                markdown.append(
                    caixa_codigo(linha)
                )

            else:

                markdown.append(
                    f"- {linha}\n"
                )

        markdown.append("\n---\n")

    salvar_md(nome, markdown)
#========================================
#SALVAR
#========================================

def salvar_md(nome, conteudo):

    if not os.path.exists(PASTA_SAIDA):
        os.makedirs(PASTA_SAIDA)

    caminho = os.path.join(
        PASTA_SAIDA,
        nome + ".md"
    )

    with open(caminho, "w", encoding="utf-8") as f:

        f.writelines(conteudo)

    print(f"✅ Gerado: {caminho}")
#========================================
#PROCESSAR PASTA
#========================================

def processar():

    arquivos = os.listdir(PASTA_ENTRADA)

    for arquivo in arquivos:

        if arquivo.endswith(".pptx"):

            caminho = os.path.join(
                PASTA_ENTRADA,
                arquivo
            )

            print(f"📄 Processando: {arquivo}")

        converter_pptx(caminho)

if __name__ == "__main__":
    processar()

print("\n🚀 Conversão concluída!")